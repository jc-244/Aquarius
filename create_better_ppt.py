import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

prs = Presentation()
# 宽屏比例
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_bg_color(slide, r, g, b):
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    rect = slide.shapes.add_shape(1, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(r, g, b)
    rect.line.fill.background()

def create_slide(title_text, content_text, bg_color=(240, 240, 240), font_color=(0,0,0)):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # 6为空白布局
    add_bg_color(slide, *bg_color)
    
    # 标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(*font_color)
    
    # 正文
    cxBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(4.5))
    cf = cxBox.text_frame
    cf.word_wrap = True
    for line in content_text.split('\n'):
        cp = cf.add_paragraph()
        cp.text = line
        cp.font.size = Pt(28)
        cp.font.color.rgb = RGBColor(*font_color)
    return slide

# 制作多页PPT
create_slide("双城记：哥伦比亚大学 vs 纽约大学", "The Tale of Two Universities in NYC\n\n极简美学版学术研究", (26, 35, 60), (255, 255, 255))
create_slide("序言：纽约的教育双子星", "在世界的中心，曼哈顿孕育了两所顶尖的高等学府。\n\n• 哥伦比亚大学：晨边高地的古典王冠\n• 纽约大学：格林威治村的无界先锋\n\n它们代表了两种截然不同的教育哲学与生活方式。", (245, 245, 245), (40, 40, 40))
create_slide("Columbia - 校园美学", "「在喧嚣的纽约，保留一块静谧的古典之地」\n\n• 标志性的封闭式传统校园 (Campus-based)\n• 壮丽的新古典主义建筑，如洛氏图书馆 (Low Library)\n• 严谨庄重的常春藤学术圣地", (225, 235, 245), (30, 45, 65))
create_slide("Columbia - 学术核心", "「博雅教育的坚实堡垒」\n\n• 招牌的 Core Curriculum（核心课程），所有人必读西方经典\n• 强势领域：新闻学（普利策奖诞生地）、法学、商科与国际关系\n• 更适合：喜欢严谨学究氛围、对经典人文底蕴有追求的精英分子", (210, 220, 235), (30, 40, 50))
create_slide("NYU - 校园精神", "「以整个纽约城市作为校园」\n\n• 完全的开放式校园 (City-based university)\n• 以华盛顿广场公园为中心，与曼哈顿繁忙的街区融为一体\n• 自由、多元、极具创造力的街头氛围", (87, 6, 140), (255, 255, 255))
create_slide("NYU - 学术态度", "「实用主义与艺术先锋的最前沿」\n\n• 极强的职业导向与实践精神，学生直接毗邻华尔街与百老汇\n• 强势领域：帝势艺术学院 (电影导演摇篮)、斯特恩商学院 (顶级投行跳板)\n• 更适合：渴望提前融入社会、崇尚实用和创造力的都市弄潮儿", (110, 30, 160), (245, 245, 245))
create_slide("核心对比矩阵", "🏫 校园体验\n• 哥大：古典的象牙塔  |  纽大：真实的城市脉动\n\n📚 学科倾向\n• 哥大：经典理论与传统学术  |  纽大：行业跳板与前沿实践\n\n👔 毕业生代表画像\n• 哥大：稳重严谨的学者与政商名流\n• 纽大：充满打破常规活力的艺术家与创业先锋", (250, 245, 240), (45, 45, 45))
create_slide("结语", "「纽约，总能容纳你所有的野心和梦想。」\n\n选择不同的学校，就是选择在纽约体验生活的一万种方式中，\n最适合你的那一种。", (40, 40, 40), (230, 230, 230))

desktop_path = os.path.expanduser("~/Desktop/哥大与纽大_重制极简版.pptx")
prs.save(desktop_path)
print("PPT done")
